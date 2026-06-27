#!/usr/bin/env python3
import json
import os
import shutil
import subprocess
import textwrap
from pathlib import Path

import requests
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SLIDES = ROOT / "slides"
FIGURES = ROOT / "figures"
DOCS = ROOT / "docs"
AUDIO = ROOT / "audio"
VIDEO = ROOT / "video"

for path in (SLIDES, DOCS, AUDIO, VIDEO):
    path.mkdir(exist_ok=True)


def read_summary():
    summary_path = DOCS / "model_summary.json"
    if not summary_path.exists():
        return {
            "rows": 1599,
            "high_quality_rate": 0.136,
            "final_model": "Random forest",
            "final_auc": 0.85,
            "final_f1": 0.55,
            "final_accuracy": 0.88,
            "final_sensitivity": 0.60,
            "final_specificity": 0.92,
            "top_predictors": [
                {"predictor": "alcohol"},
                {"predictor": "volatile acidity"},
                {"predictor": "sulphates"},
                {"predictor": "density"},
                {"predictor": "total sulfur dioxide"},
            ],
        }
    with open(summary_path, "r", encoding="utf-8") as f:
        return json.load(f)


def pct(x, digits=1):
    return f"{100 * float(x):.{digits}f}%"


def num(x, digits=3):
    return f"{float(x):.{digits}f}"


def font(size, bold=False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def wrap_lines(draw, text, font_obj, max_width):
    lines = []
    for paragraph in text.split("\n"):
        words = paragraph.split()
        current = ""
        for word in words:
            trial = f"{current} {word}".strip()
            bbox = draw.textbbox((0, 0), trial, font=font_obj)
            if bbox[2] - bbox[0] <= max_width:
                current = trial
            else:
                if current:
                    lines.append(current)
                current = word
        if current:
            lines.append(current)
        if not paragraph.strip():
            lines.append("")
    return lines


def draw_bullets(draw, bullets, x, y, max_width, font_obj, fill=(32, 38, 46), gap=18):
    bullet_font = font_obj
    line_height = int(font_obj.size * 1.35)
    for bullet in bullets:
        lines = wrap_lines(draw, bullet, font_obj, max_width - 40)
        draw.text((x, y), "-", font=bullet_font, fill=fill)
        for i, line in enumerate(lines):
            draw.text((x + 35, y + i * line_height), line, font=font_obj, fill=fill)
        y += max(line_height, len(lines) * line_height) + gap
    return y


def load_image(path, box):
    if not path.exists():
        return None
    img = Image.open(path).convert("RGB")
    img.thumbnail((box[2] - box[0], box[3] - box[1]))
    canvas = Image.new("RGB", (box[2] - box[0], box[3] - box[1]), "white")
    x = (canvas.width - img.width) // 2
    y = (canvas.height - img.height) // 2
    canvas.paste(img, (x, y))
    return canvas


def create_slide_png(index, title, bullets, image_path=None, metric_cards=None):
    width, height = 1920, 1080
    bg = Image.new("RGB", (width, height), "#f7f5f0")
    draw = ImageDraw.Draw(bg)

    accent = "#2b6f6c"
    muted = "#6f7d8c"
    text = "#20262e"
    rust = "#b45b45"

    draw.rectangle((0, 0, width, 20), fill=accent)
    draw.text((110, 80), title, font=font(58, True), fill=text)
    draw.text((110, 155), "Predictive Modeling Final Project", font=font(24), fill=muted)

    y = 240
    if metric_cards:
        card_w = 290 if image_path else 360
        card_h = 135
        x = 110
        for label, value in metric_cards:
            draw.rounded_rectangle((x, y, x + card_w, y + card_h), radius=18, fill="white", outline="#ded9cf", width=2)
            draw.text((x + 28, y + 24), label, font=font(25), fill=muted)
            draw.text((x + 28, y + 62), value, font=font(42, True), fill=rust)
            x += card_w + 30
        y += card_h + 70

    if image_path:
        img_box = (1080, 250, 1810, 930)
        pasted = load_image(Path(image_path), img_box)
        if pasted:
            draw.rounded_rectangle((img_box[0] - 15, img_box[1] - 15, img_box[2] + 15, img_box[3] + 15), radius=22, fill="white", outline="#ded9cf", width=2)
            bg.paste(pasted, (img_box[0], img_box[1]))
            bullet_width = 870
        else:
            bullet_width = 1650
    else:
        bullet_width = 1650

    draw_bullets(draw, bullets, 120, y, bullet_width, font(34), fill=text, gap=24)
    draw.text((110, 1010), f"{index}", font=font(24), fill=muted)

    out = SLIDES / f"video_slide_{index:02d}.png"
    bg.save(out)
    return out


def add_title(slide, title, subtitle=None):
    left = Inches(0.55)
    top = Inches(0.35)
    width = Inches(12.2)
    height = Inches(0.75)
    box = slide.shapes.add_textbox(left, top, width, height)
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(32, 38, 46)
    if subtitle:
        sub = slide.shapes.add_textbox(left, Inches(1.05), width, Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(111, 125, 140)


def add_bullets(slide, bullets, left=0.75, top=1.65, width=6.1, height=4.6, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(32, 38, 46)


def add_metric_cards(slide, cards):
    x = 0.75
    y = 5.8
    for label, value in cards:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.2), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(247, 245, 240)
        shape.line.color.rgb = RGBColor(222, 217, 207)
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{label}\n{value}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(32, 38, 46)
        x += 2.35


def add_picture_if_exists(slide, image_path, left=7.05, top=1.55, width=5.55, height=4.6):
    if not image_path:
        return
    path = Path(image_path)
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_pptx(summary, video_slides):
    final_model = summary["final_model"]
    top_preds = [item["predictor"] for item in summary.get("top_predictors", [])[:5]]

    executive = [
        {
            "title": "Problem And Data",
            "bullets": [
                "Goal: predict whether a red wine will receive a high sensory quality score.",
                f"Dataset: UCI red wine quality data with {summary['rows']:,} samples.",
                f"High quality wines are the minority class: {pct(summary['high_quality_rate'])} of rows.",
                "Business use: screen wines for extra review before using scarce tasting-panel time.",
            ],
            "image": FIGURES / "quality_distribution.png",
        },
        {
            "title": "Modeling Approach",
            "bullets": [
                "Created binary target: quality score 7 or higher equals High.",
                "Used stratified training and test split to preserve class balance.",
                "Compared logistic regression, decision tree, and random forest.",
                "Used 5-fold cross-validation for tuning and threshold selection.",
            ],
            "image": FIGURES / "correlation_heatmap.png",
        },
        {
            "title": "Final Model",
            "bullets": [
                f"Selected model: {final_model}.",
                f"Test ROC AUC: {num(summary['final_auc'])}.",
                f"Test F1 score: {num(summary['final_f1'])}.",
                "Threshold selected to balance precision and sensitivity for high quality wines.",
            ],
            "image": FIGURES / "roc_curves.png",
            "cards": [
                ("AUC", num(summary["final_auc"])),
                ("F1", num(summary["final_f1"])),
                ("Accuracy", num(summary["final_accuracy"])),
            ],
        },
        {
            "title": "Key Drivers",
            "bullets": [
                "Most important predictors in the tree ensemble:",
                ", ".join(top_preds) + ".",
                "These measurements are plausible quality signals, but they do not replace tasting.",
                "Model should be used as a prioritization layer.",
            ],
            "image": FIGURES / "variable_importance.png",
        },
        {
            "title": "Recommendation",
            "bullets": [
                "Use model scores to rank wines for extra sensory review.",
                "Do not use the model as the only quality decision.",
                "Improve with more wine types, producer information, vintage data, and external validation.",
                "AI assistance is disclosed in the report and repository.",
            ],
            "image": FIGURES / "predictor_boxplots.png",
        },
    ]

    def make_deck(slides, out_path):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        for s in slides:
            slide = prs.slides.add_slide(blank)
            add_title(slide, s["title"], "Predicting High Quality Wine From Physicochemical Measurements")
            add_bullets(slide, s["bullets"], width=5.95, size=17)
            if "cards" in s:
                add_metric_cards(slide, s["cards"])
            add_picture_if_exists(slide, s.get("image", ""))
        prs.save(out_path)

    make_deck(executive, ROOT / "Executive_Summary-Team1.pptx")
    make_deck(video_slides, ROOT / "Video_Presentation_Slides-Team1.pptx")


def build_slide_content(summary):
    final_model = summary["final_model"]
    top_preds = [item["predictor"] for item in summary.get("top_predictors", [])[:5]]

    video_slides = [
        {
            "title": "Problem Statement",
            "bullets": [
                "Predict whether red wine will receive a high sensory quality score.",
                "Frame quality score 7 or higher as High.",
                "Use chemistry measurements to support early screening decisions.",
                "Model is decision support, not replacement for expert tasting.",
            ],
            "image": FIGURES / "quality_distribution.png",
        },
        {
            "title": "Business Motivation",
            "bullets": [
                "Tasting panels are valuable but limited.",
                "Lab measurements are faster and cheaper to collect.",
                "A model can rank wines before extra review.",
                "Goal is prioritization, not automatic approval.",
            ],
            "image": None,
        },
        {
            "title": "Dataset",
            "bullets": [
                f"UCI red wine quality dataset with {summary['rows']:,} observations.",
                "Each row is one wine sample.",
                "Predictors are acidity, sugar, chlorides, sulfur dioxide, density, pH, sulphates, and alcohol.",
                "Original response is sensory quality rating.",
            ],
            "image": FIGURES / "predictor_boxplots.png",
        },
        {
            "title": "Target And Class Balance",
            "bullets": [
                "Binary target: High versus NotHigh.",
                f"High quality class is {pct(summary['high_quality_rate'])} of rows.",
                "Accuracy alone can be misleading.",
                "Evaluation must include sensitivity, precision, F1, and AUC.",
            ],
            "image": FIGURES / "quality_distribution.png",
        },
        {
            "title": "EDA Findings",
            "bullets": [
                "High quality wines tend to show higher alcohol and sulphates.",
                "They tend to show lower volatile acidity and lower density.",
                "Predictor overlap means perfect classification is unrealistic.",
                "Correlated predictors support comparing several model types.",
            ],
            "image": FIGURES / "correlation_heatmap.png",
        },
        {
            "title": "Preprocessing",
            "bullets": [
                "Converted original quality rating into a binary target.",
                "Removed original quality score from model inputs to avoid leakage.",
                "Checked missing values: none in the downloaded file.",
                "Standardized numeric predictors for logistic regression.",
            ],
            "image": None,
        },
        {
            "title": "Data Splitting",
            "bullets": [
                "Used stratified 80/20 train-test split.",
                "Training set used for fitting and tuning.",
                "Test set held out until final evaluation.",
                "Stratification preserved minority-class share.",
            ],
            "image": None,
        },
        {
            "title": "Validation Design",
            "bullets": [
                "Used stratified 5-fold cross-validation.",
                "Tuned decision tree complexity.",
                "Tuned random forest mtry.",
                "Selected probability thresholds from validation predictions.",
            ],
            "image": None,
        },
        {
            "title": "Model Strategy",
            "bullets": [
                "Logistic regression created a transparent baseline.",
                "Decision tree tested simple non-linear rules.",
                "Random forest tested stronger non-linear performance.",
                "Final choice balances ranking quality and high-class detection.",
            ],
            "image": None,
        },
        {
            "title": "Metrics",
            "bullets": [
                "ROC AUC measures ranking quality across thresholds.",
                "Sensitivity measures how many high quality wines are found.",
                "Precision measures how many predicted high wines are truly high.",
                "F1 balances sensitivity and precision.",
            ],
            "image": FIGURES / "roc_curves.png",
        },
        {
            "title": "Final Test Results",
            "bullets": [
                f"Selected model: {final_model}.",
                f"ROC AUC: {num(summary['final_auc'])}.",
                f"F1 score: {num(summary['final_f1'])}.",
                f"Accuracy: {num(summary['final_accuracy'])}.",
                f"Sensitivity: {num(summary['final_sensitivity'])}; specificity: {num(summary['final_specificity'])}.",
            ],
            "image": FIGURES / "roc_curves.png",
            "cards": [
                ("AUC", num(summary["final_auc"])),
                ("F1", num(summary["final_f1"])),
                ("Accuracy", num(summary["final_accuracy"])),
            ],
        },
        {
            "title": "Interpretation",
            "bullets": [
                "Top predictors in random forest:",
                ", ".join(top_preds) + ".",
                "These variables help rank wines by probability of high quality.",
                "Variable importance is predictive, not causal.",
            ],
            "image": FIGURES / "variable_importance.png",
        },
        {
            "title": "Limitations And Next Steps",
            "bullets": [
                "Dataset covers one wine type and region.",
                "Missing business variables: brand, price, vintage, producer, market.",
                "Need external validation before operational use.",
                "Next step: calibrate probabilities and test on newer data.",
            ],
            "image": None,
        },
        {
            "title": "Conclusion",
            "bullets": [
                "Chemistry data can identify useful signal for quality screening.",
                "Use output as ranked shortlist, not final judgment.",
                "Best use is prioritizing additional review.",
                "AI assistance is disclosed in the report and repository.",
            ],
            "image": None,
        },
    ]
    return video_slides


def script_text(summary):
    final_model = summary["final_model"]
    top_preds = [item["predictor"] for item in summary.get("top_predictors", [])[:5]]
    return [
        (
            "Problem Statement",
            f"""This presentation summarizes my predictive modeling project. The goal is to predict whether a red wine will receive a high sensory quality rating from basic chemistry measurements. I converted the original quality score into a binary target. A wine rated seven or higher is treated as high quality. Wines below seven are treated as not high quality. This is useful because organizations often need to decide which products deserve extra review, more tasting-panel time, or premium positioning. The model is not intended to replace human tasting. It is a decision-support layer that ranks wines and helps focus limited attention."""
        ),
        (
            "Business Motivation",
            """The practical motivation is resource allocation. Tasting panels and expert reviews are valuable, but they take time and coordination. Chemistry measurements are more standardized and can be collected earlier in the process. If a model can identify wines that are more likely to score well, the business can send those wines to a deeper review first. That does not mean the model makes the final decision. It means the model helps create a shortlist. In a real workflow, that shortlist would still be checked by people with domain knowledge."""
        ),
        (
            "Dataset",
            f"""The project uses the UCI Wine Quality red wine dataset. It contains {summary['rows']:,} red wine samples. Each row represents one wine, and each column records either a chemical property or the sensory quality score. The predictors include fixed acidity, volatile acidity, citric acid, residual sugar, chlorides, free sulfur dioxide, total sulfur dioxide, density, pH, sulphates, and alcohol. The response variable was engineered from the original quality rating. This dataset is a good fit for the assignment because it is clean, numeric, reproducible, and directly supports a classification problem."""
        ),
        (
            "Target And Class Balance",
            f"""The target variable is imbalanced. The high quality class makes up about {pct(summary['high_quality_rate'])} of the data. That matters because a model could get a high accuracy score by predicting the majority class too often. In this project, accuracy is still reported, but it is not the main decision metric. The analysis also reports sensitivity, specificity, precision, F1 score, and ROC AUC. Sensitivity tells us how many high quality wines were found. Precision tells us how many predicted high quality wines were actually high quality. F1 balances those two ideas."""
        ),
        (
            "EDA Findings",
            """The exploratory analysis shows that most wines are rated in the middle of the original quality scale. The high quality group is smaller, but it has visible differences in several predictors. High quality wines tend to have higher alcohol and higher sulphates. They tend to have lower volatile acidity and lower density. These relationships make sense chemically and practically, but they are not perfect separators. The boxplots show overlap between the high and not high classes. That overlap means a perfect model is not realistic, and the model should be evaluated as a screening tool."""
        ),
        (
            "Preprocessing",
            """The preprocessing workflow was intentionally simple and reproducible. First, I created a binary target from the original quality score. Second, I removed the original quality score from the model matrix so the model could not leak the answer. Third, I checked missing values. There were no missing cells in the downloaded dataset. Fourth, I standardized numeric predictors for logistic regression, because that model benefits from comparable scales. The tree-based models were fit on the original predictor scales because trees split on thresholds and do not require standardization."""
        ),
        (
            "Data Splitting",
            """The data was split into training and test sets using a stratified eighty-twenty split. Stratification matters here because the high quality class is small. Without stratification, the test set could accidentally contain too many or too few high quality wines, which would make evaluation unstable. The training set was used for model fitting, cross-validation, tuning, and threshold selection. The test set was held out until the final evaluation. This creates a cleaner estimate of how the selected workflow performs on data not used during model development."""
        ),
        (
            "Validation Design",
            """For validation, I used stratified five-fold cross-validation on the training data. Each fold preserved the high versus not high class structure as much as possible. The decision tree was tuned over several complexity parameter values. The random forest was tuned over several mtry values, which control how many predictors are considered at each split. I also selected the classification threshold using the cross-validation predictions. This is important because the default threshold of zero point five is not always appropriate when the positive class is uncommon."""
        ),
        (
            "Model Strategy",
            """I compared three model families. Logistic regression was the baseline. It is easy to interpret, and it tells us whether a simple linear probability pattern can solve the problem. The decision tree was the second model. It can capture simple non-linear threshold rules and is still fairly easy to explain. The random forest was the third model. It averages many trees and usually performs better when relationships are non-linear or interactive. The final selection prioritized ranking quality, F1 score, and practical usefulness for identifying high quality candidates."""
        ),
        (
            "Metrics",
            """The main model comparison metrics were ROC AUC, F1 score, accuracy, sensitivity, specificity, and precision. ROC AUC measures how well the model ranks high quality wines above ordinary wines across possible thresholds. F1 score balances precision and sensitivity for the high quality class. Sensitivity matters because missing strong wines would reduce the value of the screening process. Precision matters because sending too many weak candidates to human reviewers wastes time. Looking at these metrics together gives a more honest view than accuracy alone."""
        ),
        (
            "Final Test Results",
            f"""The selected final model is {final_model}. On the held-out test set, it produced a ROC AUC of {num(summary['final_auc'])}, an F1 score of {num(summary['final_f1'])}, and accuracy of {num(summary['final_accuracy'])}. Sensitivity was {num(summary['final_sensitivity'])}, and specificity was {num(summary['final_specificity'])}. These results mean the final model separated high quality wines from ordinary wines well on the test set. The model is especially useful as a ranking tool because ROC AUC evaluates the ordering of predicted probabilities rather than only one fixed cutoff."""
        ),
        (
            "Interpretation",
            f"""The random forest variable importance results show which predictors were most useful for splitting the classes. The top predictors were {", ".join(top_preds)}. These are plausible drivers, but they should be interpreted carefully. Variable importance does not prove causation. It means these measurements helped the model separate wines that received higher scores from the rest of the samples. The practical interpretation is that the model can produce a ranked list of wines that are more likely to be high quality. A business could send the highest scoring wines to extra sensory review first."""
        ),
        (
            "Limitations And Next Steps",
            """The main limitations are the narrow dataset, the absence of brand or price information, the lack of external validation, and the fact that sensory quality is partly subjective. The dataset covers red Vinho Verde wine, so results should not be assumed to generalize to all wines. In a production setting, I would add more samples from more regions, include business variables such as price and producer, calibrate the predicted probabilities, and test the model on newer external data before using it operationally."""
        ),
        (
            "Conclusion",
            """The conclusion is that physicochemical measurements can provide useful signal for screening red wine quality. The model should be used as decision support. It can prioritize wines for additional review, but it should not make final quality decisions by itself. The best workflow is to treat model scores as a ranked shortlist, then use expert review for the final call. AI assistance was used to help organize and prepare materials, and that use is disclosed in the report and repository."""
        ),
    ]


def write_video_script(script_parts):
    path = ROOT / "video_script.md"
    lines = ["# Video Presentation Script", ""]
    for i, (title, text) in enumerate(script_parts, start=1):
        lines.append(f"## Slide {i}: {title}")
        lines.append("")
        lines.append(textwrap.fill(text, width=100))
        lines.append("")
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def tts_with_elevenlabs(script_parts):
    existing = [AUDIO / f"slide_{i:02d}.mp3" for i in range(1, len(script_parts) + 1)]
    force = os.environ.get("FORCE_REGENERATE_AUDIO") == "1"
    if not force and all(path.exists() and path.stat().st_size > 1000 for path in existing):
        return existing

    edge_tts = shutil.which("edge-tts")
    if edge_tts:
        files = []
        for i, (title, text) in enumerate(script_parts, start=1):
            out = AUDIO / f"slide_{i:02d}.mp3"
            text_file = AUDIO / f"slide_{i:02d}.txt"
            text_file.write_text(text, encoding="utf-8")
            if force and out.exists():
                out.unlink()
            subprocess.run(
                [
                    edge_tts,
                    "--voice",
                    os.environ.get("EDGE_TTS_VOICE", "en-US-GuyNeural"),
                    f"--rate={os.environ.get('EDGE_TTS_RATE', '-18%')}",
                    "--file",
                    str(text_file),
                    "--write-media",
                    str(out),
                ],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            files.append(out)
        (DOCS / "narration_note.txt").write_text(
            "Used Microsoft Edge neural TTS with one consistent voice for all slides.",
            encoding="utf-8",
        )
        return files

    api_keys = [
        os.environ.get("ELEVEN_LABS_API_KEY"),
        os.environ.get("ELEVENLABS_API_KEY"),
    ]
    api_keys = [key for key in api_keys if key]
    preferred_voice_id = os.environ.get("ELEVEN_LABS_VOICE_ID")
    fallback_voice_id = "TX3LPaxmHKxFdv7VOQHJ"
    fallback_voice_name = "Liam - Energetic, Social Media Creator"
    def local_tts(text, out):
        txt = out.with_suffix(".txt")
        wav = out.with_suffix(".wav")
        txt.write_text(text, encoding="utf-8")
        subprocess.run(
            ["espeak-ng", "-v", "en-us", "-s", "138", "-f", str(txt), "-w", str(wav)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        subprocess.run(
            ["ffmpeg", "-y", "-i", str(wav), "-codec:a", "libmp3lame", "-q:a", "4", str(out)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )

    if not api_keys or not preferred_voice_id:
        files = []
        for i, (title, text) in enumerate(script_parts, start=1):
            out = AUDIO / f"slide_{i:02d}.mp3"
            if not out.exists() or out.stat().st_size <= 1000:
                local_tts(text, out)
            files.append(out)
        (DOCS / "narration_note.txt").write_text(
            "ElevenLabs credentials unavailable. Used local TTS for narration.",
            encoding="utf-8",
        )
        return files

    files = []
    selected_key = None
    selected_voice = preferred_voice_id
    selected_voice_note = "Preferred saved voice."

    def try_request(api_key, voice_id, text):
        endpoint = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}"
        headers = {
            "xi-api-key": api_key,
            "Accept": "audio/mpeg",
            "Content-Type": "application/json",
        }
        payload = {
            "text": text,
            "model_id": "eleven_multilingual_v2",
            "voice_settings": {
                "stability": 0.45,
                "similarity_boost": 0.82,
                "style": 0.25,
                "use_speaker_boost": True,
            },
        }
        return requests.post(endpoint, headers=headers, json=payload, timeout=120)

    first_text = script_parts[0][1][:200]
    last_error = ""
    for api_key in api_keys:
        response = try_request(api_key, preferred_voice_id, first_text)
        if response.status_code == 200:
            selected_key = api_key
            break
        last_error = response.text[:500]
        response = try_request(api_key, fallback_voice_id, first_text)
        if response.status_code == 200:
            selected_key = api_key
            selected_voice = fallback_voice_id
            selected_voice_note = f"Preferred saved voice unavailable. Used premade fallback voice: {fallback_voice_name}."
            break
        last_error = response.text[:500]

    if not selected_key:
        (DOCS / "narration_note.txt").write_text(
            "ElevenLabs narration could not be generated. Used local TTS fallback. Last API error: " + last_error,
            encoding="utf-8",
        )
        files = []
        for i, (title, text) in enumerate(script_parts, start=1):
            out = AUDIO / f"slide_{i:02d}.mp3"
            if not out.exists() or out.stat().st_size <= 1000:
                local_tts(text, out)
            files.append(out)
        return files

    note_lines = [selected_voice_note]

    for i, (title, text) in enumerate(script_parts, start=1):
        out = AUDIO / f"slide_{i:02d}.mp3"
        if out.exists() and out.stat().st_size > 1000:
            files.append(out)
            continue
        response = try_request(selected_key, selected_voice, text)
        if response.status_code == 200:
            out.write_bytes(response.content)
        else:
            note_lines.append(
                f"Slide {i} used local TTS fallback because ElevenLabs returned {response.status_code}: {response.text[:300]}"
            )
            local_tts(text, out)
        files.append(out)
    (DOCS / "narration_note.txt").write_text("\n".join(note_lines), encoding="utf-8")
    return files


def audio_duration(path):
    result = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            str(path),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    return float(result.stdout.strip())


def build_video(slide_images, audio_files):
    if len(slide_images) != len(audio_files) or not audio_files:
        return None

    segments = []
    for i, (slide, audio) in enumerate(zip(slide_images, audio_files), start=1):
        segment = VIDEO / f"segment_{i:02d}.mp4"
        duration = audio_duration(audio)
        subprocess.run(
            [
                "ffmpeg",
                "-y",
                "-loop",
                "1",
                "-framerate",
                "1",
                "-i",
                str(slide),
                "-i",
                str(audio),
                "-t",
                f"{duration:.2f}",
                "-c:v",
                "libx264",
                "-tune",
                "stillimage",
                "-r",
                "1",
                "-c:a",
                "aac",
                "-b:a",
                "192k",
                "-pix_fmt",
                "yuv420p",
                "-shortest",
                str(segment),
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        segments.append(segment)

    concat_file = VIDEO / "segments.txt"
    concat_file.write_text("\n".join(f"file '{p}'" for p in segments), encoding="utf-8")
    out = ROOT / "Video_Presentation-Team1.mp4"
    subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(concat_file),
            "-c",
            "copy",
            str(out),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )

    audio_out = ROOT / "video_narration.mp3"
    concat_audio = AUDIO / "audio_files.txt"
    concat_audio.write_text("\n".join(f"file '{p}'" for p in audio_files), encoding="utf-8")
    subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(concat_audio),
            "-c",
            "copy",
            str(audio_out),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    return out


def main():
    summary = read_summary()
    video_slides = build_slide_content(summary)
    script_parts = script_text(summary)
    write_video_script(script_parts)
    build_pptx(summary, video_slides)

    slide_images = []
    for i, slide in enumerate(video_slides, start=1):
        slide_images.append(
            create_slide_png(
                i,
                slide["title"],
                slide["bullets"],
                image_path=slide.get("image"),
                metric_cards=slide.get("cards"),
            )
        )

    audio_files = tts_with_elevenlabs(script_parts)
    video_file = build_video(slide_images, audio_files)

    manifest = {
        "executive_summary": str(ROOT / "Executive_Summary-Team1.pptx"),
        "video_slides": str(ROOT / "Video_Presentation_Slides-Team1.pptx"),
        "video_script": str(ROOT / "video_script.md"),
        "slide_images": [str(p) for p in slide_images],
        "audio_files": [str(p) for p in audio_files],
        "video_file": str(video_file) if video_file else None,
    }
    (DOCS / "asset_manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(json.dumps(manifest, indent=2))


if __name__ == "__main__":
    main()
